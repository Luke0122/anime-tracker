# -*- coding: utf-8 -*-
"""PPT 生成共享工具：16:9 画布、中文字体、表格、卡片等样式助手。"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

FONT = "微软雅黑"
ACCENT = "8A9A7B"       # 莫兰迪鼠尾草绿
HEADER_BG = "6E7D62"    # 表头深鼠尾草
LIGHT_BG = "FBFAF6"     # 卡片/斑马纹米白
BG = "F4F1EA"           # 幻灯片背景米灰
DARK = "5C564E"         # 正文深灰棕
GRAY = "8D877C"         # 次要文字灰
SCORE_HIGH = "6B8E6B"
SCORE_MID = "C9A227"
SCORE_LOW = "9A9A9A"

SLIDE_W = 13.333
SLIDE_H = 7.5


def init_prs():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string(BG)
    return slide


def set_run(run, size=10, bold=False, color=DARK, name=FONT, italic=False):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = RGBColor.from_string(color)
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def add_text(slide, x, y, w, h, text, size=12, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line_spacing=1.0,
             wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = valign
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    text = clean_text(text)
    lines = text.split("\n") if "\n" in text else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        set_run(run, size=size, bold=bold, color=color)
    return box


def add_card(slide, x, y, w, h, number, label, number_size=34, label_size=14,
             fill=LIGHT_BG, number_color=ACCENT, label_color=GRAY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = str(number)
    set_run(r1, size=number_size, bold=True, color=number_color)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(2)
    r2 = p2.add_run()
    r2.text = str(label)
    set_run(r2, size=label_size, color=label_color)
    return shape


def shade_cell(cell, color):
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor.from_string(color)


def cell_write(cell, text, size=10, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = clean_text(text).replace("\n", " ")
    set_run(run, size=size, bold=bold, color=color)


def make_table(slide, x, y, w, headers, rows, font_size=10, col_widths=None,
               row_height=0.3, header_fill=HEADER_BG, zebra=True,
               score_col=None, first_col_bold=False, row_heights=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(row_height * n_rows))
    table = shape.table
    table.first_row = False
    table.horz_banding = False
    if col_widths:
        total = sum(col_widths)
        for j, cw in enumerate(col_widths):
            table.columns[j].width = Inches(w * cw / total)
    for j, h in enumerate(headers):
        cell_write(table.cell(0, j), h, size=font_size, bold=True, color="FFFFFF",
                   align=PP_ALIGN.CENTER)
        shade_cell(table.cell(0, j), header_fill)
    table.rows[0].height = Inches(row_height)
    for i, row in enumerate(rows, start=1):
        rh = row_heights[i - 1] if row_heights else row_height
        table.rows[i].height = Inches(rh)
        for j, val in enumerate(row):
            align = PP_ALIGN.CENTER if j in (0, 2, 3, 4) else PP_ALIGN.LEFT
            cell_write(table.cell(i, j), val, size=font_size,
                       bold=(first_col_bold and j == 0), align=align)
            if score_col is not None and j == score_col:
                try:
                    v = float(str(val))
                except (TypeError, ValueError):
                    v = None
                if v is not None:
                    color = SCORE_HIGH if v >= 8.5 else (SCORE_MID if v >= 7.5 else SCORE_LOW)
                    table.cell(i, j).text_frame.paragraphs[0].runs[0].font.color.rgb = \
                        RGBColor.from_string(color)
        if zebra and i % 2 == 0:
            for j in range(n_cols):
                shade_cell(table.cell(i, j), LIGHT_BG)
    return table


def add_cover_slide(prs, title, subtitle_lines, date_line=None):
    slide = blank_slide(prs)
    add_text(slide, 1.0, 2.05, 11.333, 1.2, title, size=46, bold=True,
             color=ACCENT, align=PP_ALIGN.CENTER)
    y = 3.45
    for line in subtitle_lines:
        add_text(slide, 1.0, y, 11.333, 0.5, line, size=17, color=GRAY,
                 align=PP_ALIGN.CENTER)
        y += 0.42
    if date_line:
        add_text(slide, 1.0, y + 0.35, 11.333, 0.45, date_line, size=15,
                 color="8A8F98", align=PP_ALIGN.CENTER)
    add_text(slide, 1.0, 6.85, 11.333, 0.4, "Anime Auto · 自动生成",
             size=12, color="6B7280", align=PP_ALIGN.CENTER)
    return slide


def add_footer(slide, page_no, note="数据来源：yuc.wiki + Bangumi"):
    add_text(slide, 0.5, 7.13, 9.0, 0.28, note, size=9.5, color="8A8F98")
    add_text(slide, 11.9, 7.13, 0.933, 0.28, "第 %d 页" % page_no, size=9,
             color="8A8F98", align=PP_ALIGN.RIGHT)


def clean_text(s):
    """清理文本中的换行/控制字符，避免写入 PPT XML 后变成 _x000D_/_x000B_ 乱码。"""
    if s is None:
        return ""
    s = str(s)
    s = s.replace("\r\n", "\n").replace("\r", "\n").replace("\x0b", "\n")
    while "\n\n\n" in s:
        s = s.replace("\n\n\n", "\n\n")
    return s


def truncate(s, limit=120):
    s = "" if s is None else str(s)
    s = s.strip()
    if len(s) <= limit:
        return s
    return s[:limit].rstrip() + "……"
